import json
import os
from pathlib import Path
from typing import Optional

import fitz
from docx import Document
import openpyxl
import csv
import io


def parse_document(file_path: str) -> Optional[str]:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext in (".docx", ".doc"):
        return _parse_docx(file_path)
    elif ext in (".xlsx", ".xls"):
        return _parse_xlsx(file_path)
    elif ext == ".csv":
        return _parse_csv(file_path)
    elif ext in (".txt", ".md"):
        return _parse_text(file_path)
    elif ext == ".pptx":
        return _parse_pptx(file_path)
    return None


def _parse_pdf(file_path: str) -> str:
    text_parts = []
    with fitz.open(file_path) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n\n".join(text_parts)


def _parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paras)


def _parse_xlsx(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"=== Sheet: {sheet} ===")
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) if v is not None else "" for v in row]
            lines.append(" | ".join(vals))
    wb.close()
    return "\n".join(lines)


def _parse_csv(file_path: str) -> str:
    lines = []
    with open(file_path, newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        for row in reader:
            lines.append(" | ".join(row))
    return "\n".join(lines)


def _parse_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n\n".join(lines)


def parse_document_raw(content: bytes, ext: str) -> Optional[str]:
    import tempfile
    suffix = ext if ext.startswith(".") else "." + ext
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(content)
        tmp_path = tmp.name
    try:
        return parse_document(tmp_path)
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


def extract_metadata(text: str, filename: str) -> dict:
    import re
    metadata = {
        "filename": filename,
        "page_count_estimate": max(1, len(text) // 3000),
        "word_count": len(text.split()),
        "char_count": len(text),
        "has_dates": bool(re.search(r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b', text)),
        "has_amounts": bool(re.search(r'(?:Rs\.?|INR|USD|EUR|₹|\$)\s*[\d,]+\.?\d{0,2}', text)),
    }
    return metadata
